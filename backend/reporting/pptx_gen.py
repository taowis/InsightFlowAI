from pptx import Presentation

def build_pptx(narrative: str, out_path: str):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Monthly Performance Summary"
    slide.placeholders[1].text = narrative
    prs.save(out_path)
    return out_path
